import streamlit as st

st.set_page_config(
        page_title="EduVision",
        page_icon="📒",
)

st.title('Instructions')

import streamlit as st
import base64
import pptx
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import streamlit as st
import openai

# Load your OpenAI API key (replace 'YOUR_API_KEY' with your actual API key)
openai.api_key = st.secrets["OPENAI_API_KEY"]

# Define custom formatting options
TITLE_FONT_SIZE = Pt(32)
SLIDE_FONT_SIZE = Pt(20)
TEXT_BOX_WIDTH = Inches(8)
TEXT_BOX_HEIGHT = Inches(4.5)
SPACE_BETWEEN_SLIDES = Inches(0.5)

# Custom colors
BACKGROUND_COLOR = RGBColor(0, 160, 186)
TITLE_TEXT_COLOR = RGBColor(255, 255, 255)
SLIDE_TEXT_COLOR = RGBColor(227, 226, 228)

def generate_slide_titles(topic, num_slides=5):
    slide_titles = []
    for i in range(num_slides):
        user_message = f"Generate a title for slide {i + 1} on the topic: {topic}"
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a chatbot that generates presentation that looks attractive."},
                {"role": "user", "content": user_message},
            ],
            max_tokens=30
        )
        slide_titles.append(response.choices[0].message.content.strip())
    return slide_titles

def generate_slide_content(slide_title):
    user_message = f"Generate content for the slide titled: {slide_title}"
    response = openai.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a chatbot that generates attractive presentation and beautifull slide content."},
            {"role": "user", "content": user_message},
        ],
        max_tokens=140
    )
    return response.choices[0].message.content.strip()  

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()

    # Create a title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_TEXT_COLOR
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Set the background color of the title slide
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR

    # Create content slides
    slide_layout = prs.slide_layouts[5]
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = SLIDE_TEXT_COLOR
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set the background color of content slides
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        left_inch = Inches(0.5)
        top_inch = Inches(1.0)
        content = slide.shapes.add_textbox(left_inch, top_inch, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
        text_frame = content.text_frame
        p = text_frame.add_paragraph()
        p.text = slide_content
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        text_frame.word_wrap = True
        text_frame.paragraphs[0].font.size = SLIDE_FONT_SIZE
        text_frame.paragraphs[0].font.color.rgb = SLIDE_TEXT_COLOR

    ppt_filename = f"{topic}_presentation.pptx"
    prs.save(ppt_filename)
    return ppt_filename

def main():
    st.title("PowerPoint Presentation Generator App")

    topic = st.text_input("Enter the topic for your presentation:")
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        with st.spinner("Generating presentation... Please wait."):
            slide_titles = generate_slide_titles(topic)
            slide_contents = [generate_slide_content(title) for title in slide_titles]
            ppt_filename = create_presentation(topic, slide_titles, slide_contents)

        st.success("Presentation generated successfully!")

        # Provide a download link for the generated PowerPoint presentation
        st.markdown(get_ppt_download_link(ppt_filename), unsafe_allow_html=True)

def get_ppt_download_link(ppt_filename):
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()